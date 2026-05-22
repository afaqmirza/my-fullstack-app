import sys
import json
import os
import shutil
import zipfile
import re
import traceback
import io
from urllib.parse import quote_plus
import time
import uuid

# PDF / Doc tools
import pdfplumber
import pandas as pd
from pptx import Presentation
from pdf2docx import Converter
from pypdf import PdfWriter, PdfReader
from PIL import Image, ImageOps
import img2pdf
import fitz  # PyMuPDF
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import PyPDF2

# Intelligence / Automation
from bs4 import BeautifulSoup
try:
    from selenium import webdriver
    from selenium.webdriver.chrome.options import Options
    from selenium.webdriver.common.by import By
except ImportError:
    webdriver = Options = By = None
from groq import Groq
from docx import Document as DocxDocument
from deep_translator import GoogleTranslator
from docx_utils import extract_docx_text_safe, extract_docx_tables_safe
from tesseract_ocr import ocr_file, write_text_output, tesseract_status

try:
    import pythoncom
    import win32com.client as win32
except ImportError:
    pythoncom = None
    win32 = None

SOCIAL_PATTERNS = {
    "facebook":  re.compile(r"https?://(www\.)?facebook\.com/(?!sharer|share|dialog|plugins|tr\b)[^\s\"'<>?&]+"),
    "instagram": re.compile(r"https?://(www\.)?instagram\.com/[^\s\"'<>?&]+"),
    "twitter":   re.compile(r"https?://(www\.)?(twitter\.com|x\.com)/[^\s\"'<>?&]+"),
    "linkedin":  re.compile(r"https?://(www\.)?linkedin\.com/(company|in)/[^\s\"'<>?&]+"),
    "youtube":   re.compile(r"https?://(www\.)?youtube\.com/(channel|c|user|@)[^\s\"'<>?&]+"),
}

def extract_social_links(text: str) -> dict:
    result = {}
    for platform, pat in SOCIAL_PATTERNS.items():
        matches = [m.group(0) for m in pat.finditer(text)]
        seen = set()
        unique = []
        for m in matches:
            base = m.rstrip("/")
            if base not in seen:
                seen.add(base)
                unique.append(base)
        if unique:
            result[platform] = unique[0]
    return result

def parse_page_ranges(range_str, total_pages):
    pages = set()
    parts = [p.strip() for p in range_str.split(",") if p.strip()]
    for part in parts:
        if "-" in part:
            try:
                start, end = part.split("-", 1)
                start = int(start.strip())
                end = int(end.strip())
                for p in range(start, end + 1):
                    if 1 <= p <= total_pages:
                        pages.add(p - 1)
            except: continue
        else:
            try:
                p = int(part.strip())
                if 1 <= p <= total_pages:
                    pages.add(p - 1)
            except: continue
    return sorted(pages)

def process_task(task, files, args, temp_dir):
    os.makedirs(temp_dir, exist_ok=True)
    
    if task == "word-to-pdf":
        input_path = files[0]
        output_path = os.path.join(temp_dir, f"{uuid.uuid4()}_output.pdf")
        if pythoncom is None or win32 is None:
            raise Exception("pywin32 is required for Word conversion on Windows.")
        pythoncom.CoInitialize()
        word = win32.Dispatch("Word.Application")
        word.Visible = False
        word.DisplayAlerts = 0
        doc = word.Documents.Open(os.path.abspath(input_path))
        doc.SaveAs(os.path.abspath(output_path), FileFormat=17)
        doc.Close(False)
        word.Quit()
        pythoncom.CoUninitialize()
        return output_path

    elif task == "word-to-excel":
        input_path = files[0]
        output_path = os.path.join(temp_dir, f"{uuid.uuid4()}_output.xlsx")
        
        is_doc = input_path.lower().endswith(".doc")
        docx_path = input_path
        
        if is_doc:
            if pythoncom is None or win32 is None:
                raise Exception("pywin32 is required to process legacy .doc files on Windows.")
            docx_path = os.path.join(temp_dir, f"{uuid.uuid4()}_converted.docx")
            pythoncom.CoInitialize()
            word = win32.Dispatch("Word.Application")
            word.Visible = False
            word.DisplayAlerts = 0
            try:
                doc = word.Documents.Open(os.path.abspath(input_path))
                doc.SaveAs(os.path.abspath(docx_path), FileFormat=16) # wdFormatDocumentDefault
                doc.Close(False)
            finally:
                word.Quit()
                pythoncom.CoUninitialize()
        
        docx_text = extract_docx_text_safe(docx_path)
        docx_tables = extract_docx_tables_safe(docx_path)

        with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
            # 1. Process paragraphs
            sentences = []
            if docx_text.strip():
                for txt in docx_text.split("\n"):
                    txt = txt.strip()
                    if txt:
                        s_list = re.split(r'(?<=[.!?])\s+(?=[A-Z])', txt)
                        sentences.extend([s.strip() for s in s_list if s.strip()])
            
            if sentences:
                df_txt = pd.DataFrame(sentences, columns=["Document Content"])
                df_txt.to_excel(writer, sheet_name="Text Content", index=False)
                ws = writer.sheets["Text Content"]
                ws.column_dimensions["A"].width = 120
            
            # 2. Process tables
            for idx, table_data in enumerate(docx_tables):
                if table_data:
                    df_table = pd.DataFrame(table_data)
                    sheet_name = f"Table {idx + 1}"[:31]
                    df_table.to_excel(writer, sheet_name=sheet_name, index=False, header=False)
                    
                    # Auto-fit columns
                    ws = writer.sheets[sheet_name]
                    for col in ws.columns:
                        vals = [str(cell.value or '') for cell in col]
                        max_len = max(len(v) for v in vals) if vals else 0
                        col_letter = col[0].column_letter
                        ws.column_dimensions[col_letter].width = max(max_len + 3, 10)
                        
            # If nothing was written, write an empty note
            if not sentences and not docx_tables:
                df_empty = pd.DataFrame(["No text or tables found in the Word document."], columns=["Notification"])
                df_empty.to_excel(writer, sheet_name="Sheet1", index=False)
                
        return output_path

    elif task == "pdf-to-ppt":
        input_path = files[0]
        output_path = os.path.join(temp_dir, f"{uuid.uuid4()}_output.pptx")
        prs = Presentation()
        with pdfplumber.open(input_path) as pdf:
            for page in pdf.pages:
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
                img_path = os.path.join(temp_dir, f"page_{page.page_number}.png")
                page.to_image(resolution=150).save(img_path)
                slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
                os.remove(img_path)
        prs.save(output_path)
        return output_path

    elif task == "pdf-to-image":
        input_path = files[0]
        output_path = os.path.join(temp_dir, f"{uuid.uuid4()}_images.zip")
        image_dir = os.path.join(temp_dir, f"{uuid.uuid4()}_images_dir")
        os.makedirs(image_dir, exist_ok=True)
        with pdfplumber.open(input_path) as pdf:
            for page in pdf.pages:
                image_path = os.path.join(image_dir, f"page_{page.page_number}.png")
                page.to_image(resolution=300).save(image_path)
        with zipfile.ZipFile(output_path, 'w') as zipf:
            for root, _, fs in os.walk(image_dir):
                for f in fs:
                    zipf.write(os.path.join(root, f), f)
        return output_path

    elif task == "pdf-to-excel":
        input_path = files[0]
        output_path = os.path.join(temp_dir, f"{uuid.uuid4()}_output.xlsx")
        all_text = ""
        with pdfplumber.open(input_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    all_text += " " + text
        clean_text = " ".join(all_text.split())
        sentences = re.split(r'(?<=[.!?])\s+(?=[A-Z])', clean_text)
        sentences = [s.strip() for s in sentences if s.strip()]
        if sentences:
            df = pd.DataFrame(sentences, columns=["Sentence"])
            with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
                df.to_excel(writer, index=False)
                ws = writer.sheets["Sheet1"]
                ws.column_dimensions["A"].width = 120
        else:
            raise Exception("No content found in PDF to convert.")
        return output_path

    elif task == "pdf-to-word":
        input_path = files[0]
        output_path = os.path.join(temp_dir, f"{uuid.uuid4()}_output.docx")
        cv = Converter(input_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()
        return output_path

    elif task == "pdf-merge":
        output_path = os.path.join(temp_dir, f"{uuid.uuid4()}_output.pdf")
        writer = PdfWriter()
        for input_path in files:
            reader = PdfReader(input_path)
            for page in reader.pages:
                writer.add_page(page)
        with open(output_path, "wb") as out_file:
            writer.write(out_file)
        return output_path

    elif task == "pdf-split":
        input_path = files[0]
        mode = args.get("mode", "range")
        reader = PdfReader(input_path)
        total_pages = len(reader.pages)
        
        if mode == "range":
            range_str = args.get("range_str", "")
            pages = parse_page_ranges(range_str, total_pages)
            if not pages:
                raise Exception("No valid pages found in range.")
            writer = PdfWriter()
            for idx in pages:
                writer.add_page(reader.pages[idx])
            output_path = os.path.join(temp_dir, f"{uuid.uuid4()}_output.pdf")
            with open(output_path, "wb") as f:
                writer.write(f)
            return output_path
        else:
            output_path = os.path.join(temp_dir, f"{uuid.uuid4()}_split.zip")
            image_dir = os.path.join(temp_dir, f"{uuid.uuid4()}_split_dir")
            os.makedirs(image_dir, exist_ok=True)
            
            if mode == "each":
                for i in range(total_pages):
                    writer = PdfWriter()
                    writer.add_page(reader.pages[i])
                    out_part = os.path.join(image_dir, f"page_{i+1}.pdf")
                    with open(out_part, "wb") as f:
                        writer.write(f)
            elif mode == "parts":
                parts_count = int(args.get("parts_count", 2))
                pages_per_part = total_pages // parts_count
                remainder = total_pages % parts_count
                start = 0
                for part in range(parts_count):
                    extra = 1 if part < remainder else 0
                    end = start + pages_per_part + extra
                    if start == end: continue
                    writer = PdfWriter()
                    for idx in range(start, end):
                        writer.add_page(reader.pages[idx])
                    out_part = os.path.join(image_dir, f"part_{part+1}.pdf")
                    with open(out_part, "wb") as f:
                        writer.write(f)
                    start = end
            with zipfile.ZipFile(output_path, 'w') as zipf:
                for root, _, fs in os.walk(image_dir):
                    for f in fs:
                        zipf.write(os.path.join(root, f), f)
            return output_path

    elif task == "pdf-compress":
        input_path = files[0]
        output_path = os.path.join(temp_dir, f"{uuid.uuid4()}_output.pdf")
        reader = PdfReader(input_path)
        writer = PdfWriter()
        for page in reader.pages:
            page.compress_content_streams()
            writer.add_page(page)
        writer.compress_identical_objects(remove_identicals=True, remove_orphans=True)
        with open(output_path, "wb") as f:
            writer.write(f)
        return output_path

    elif task == "image-to-pdf":
        output_path = os.path.join(temp_dir, f"{uuid.uuid4()}_output.pdf")
        image_paths = []
        for file_path in files:
            img = Image.open(file_path)
            img = ImageOps.exif_transpose(img)
            if img.mode in ("RGBA", "LA", "P"):
                bg = Image.new("RGB", img.size, (255, 255, 255))
                if img.mode == "P":
                    img = img.convert("RGBA")
                bg.paste(img, mask=img.split()[-1] if img.mode == "RGBA" else None)
                img = bg
            elif img.mode != "RGB":
                img = img.convert("RGB")
            buf_path = os.path.join(temp_dir, f"processed_{uuid.uuid4()}.jpg")
            img.save(buf_path, format="JPEG", quality=95)
            image_paths.append(buf_path)
        if not image_paths:
            raise Exception("No valid images processed.")
        pdf_bytes = img2pdf.convert(image_paths)
        with open(output_path, "wb") as f:
            f.write(pdf_bytes)
        return output_path

    elif task == "pdf-protect":
        input_path = files[0]
        password = args.get("password")
        output_path = os.path.join(temp_dir, f"{uuid.uuid4()}_output.pdf")
        reader = PdfReader(input_path)
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        writer.encrypt(user_password=password, owner_password=password, use_128bit=True)
        with open(output_path, "wb") as f:
            writer.write(f)
        return output_path

    elif task == "pdf-unlock":
        input_path = files[0]
        password = args.get("password")
        output_path = os.path.join(temp_dir, f"{uuid.uuid4()}_output.pdf")
        reader = PdfReader(input_path)
        if reader.is_encrypted:
            if reader.decrypt(password) == 0:
                raise Exception("Incorrect password.")
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        with open(output_path, "wb") as f:
            writer.write(f)
        return output_path

    elif task == "pdf-rotate":
        input_path = files[0]
        angle = int(args.get("angle", 90))
        pages_str = args.get("pages", "ALL")
        output_path = os.path.join(temp_dir, f"{uuid.uuid4()}_output.pdf")
        reader = PdfReader(input_path)
        writer = PdfWriter()
        total_pages = len(reader.pages)
        target_pages = parse_page_ranges(pages_str, total_pages) if pages_str.upper() != "ALL" else range(total_pages)
        for i, page in enumerate(reader.pages):
            if i in target_pages:
                page.rotate(angle)
            writer.add_page(page)
        with open(output_path, "wb") as f:
            writer.write(f)
        return output_path

    elif task == "pdf-watermark":
        input_path = files[0]
        text = args.get("text", "Watermark")
        opacity = float(args.get("opacity", 0.3))
        output_path = os.path.join(temp_dir, f"{uuid.uuid4()}_output.pdf")
        
        packet = io.BytesIO()
        can = canvas.Canvas(packet, pagesize=letter)
        can.setFillColorRGB(0, 0, 0, alpha=opacity)
        can.setFont("Helvetica", 40)
        can.saveState()
        can.translate(200, 400)
        can.rotate(45)
        can.drawString(0, 0, text)
        can.restoreState()
        can.save()
        packet.seek(0)
        
        watermark_reader = PdfReader(packet)
        wm_page = watermark_reader.pages[0]
        
        reader = PdfReader(input_path)
        writer = PdfWriter()
        for page in reader.pages:
            page.merge_page(wm_page)
            writer.add_page(page)
        with open(output_path, "wb") as f:
            writer.write(f)
        return output_path

    elif task == "pdf-edit":
        input_path = files[0]
        edit_text = args.get("edit_text")
        x = int(args.get("x", 100))
        y = int(args.get("y", 100))
        output_path = os.path.join(temp_dir, f"{uuid.uuid4()}_output.pdf")
        doc = fitz.open(input_path)
        page = doc[0]
        if edit_text:
            page.insert_text((x, y), edit_text, fontsize=14, color=(0, 0, 0))
        if len(files) > 1: # second file is image
            img_path = files[1]
            rect = fitz.Rect(x, y, x + 150, y + 150)
            page.insert_image(rect, filename=img_path)
        doc.save(output_path)
        return output_path

    elif task == "html-to-docx":
        import base64
        from docx.shared import Inches, Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        html_content = ""
        if files and len(files) > 0:
            with open(files[0], "r", encoding="utf-8", errors="ignore") as f:
                html_content = f.read()
        else:
            html_content = args.get("html", "")

        output_path = os.path.join(temp_dir, f"{uuid.uuid4()}_output.docx")
        doc = DocxDocument()

        def add_image_from_src(src, width_px=None):
            if not src or not src.startswith("data:"):
                return
            try:
                header, b64data = src.split(",", 1)
                raw = base64.b64decode(b64data)
                stream = io.BytesIO(raw)
                width = Inches(2.5)
                if width_px:
                    width = Inches(max(0.5, width_px / 96))
                doc.add_picture(stream, width=width)
            except Exception:
                pass

        def add_paragraph_from_element(el, bold=False, italic=False, size_pt=None):
            text = el.get_text(strip=True) if hasattr(el, "get_text") else str(el).strip()
            if not text:
                return
            p = doc.add_paragraph()
            run = p.add_run(text)
            run.bold = bold
            run.italic = italic
            if size_pt:
                run.font.size = Pt(size_pt)

        soup = BeautifulSoup(html_content or "<p></p>", "html.parser")

        for wm in soup.find_all(class_="doc-watermark"):
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = p.add_run(wm.get_text(strip=True))
            run.font.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
            size = 48
            style = wm.get("style", "")
            m = re.search(r"font-size:\s*(\d+)", style)
            if m:
                size = max(24, min(96, int(m.group(1)) // 2))
            run.font.size = Pt(size)
            run.bold = True

        body_nodes = soup.find_all(class_="doc-body")
        root = body_nodes[0] if body_nodes else soup

        def walk(node):
            if getattr(node, "name", None) is None:
                return
            if node.name == "img":
                add_image_from_src(node.get("src", ""))
                return
            if node.name in ("h1", "h2", "h3"):
                add_paragraph_from_element(node, bold=True, size_pt=16 if node.name == "h1" else 14)
                return
            if node.name == "p":
                imgs = node.find_all("img")
                if imgs:
                    for img in imgs:
                        w = img.get("style", "")
                        wm = re.search(r"width:\s*(\d+)", w)
                        add_image_from_src(img.get("src", ""), int(wm.group(1)) if wm else None)
                if node.get_text(strip=True):
                    add_paragraph_from_element(node)
                return
            if node.name in ("ul", "ol"):
                for li in node.find_all("li", recursive=False):
                    add_paragraph_from_element(li)
                return
            if node.name == "table":
                rows = node.find_all("tr")
                if rows:
                    cols = max(len(r.find_all(["td", "th"])) for r in rows)
                    table = doc.add_table(rows=len(rows), cols=cols)
                    for ri, tr in enumerate(rows):
                        cells = tr.find_all(["td", "th"])
                        for ci, cell in enumerate(cells):
                            if ci < cols:
                                table.rows[ri].cells[ci].text = cell.get_text(strip=True)
                return
            if node.get("class") and "editor-overlay-block" in node.get("class", []):
                img = node.find("img")
                if img:
                    w = img.get("style", "")
                    wm = re.search(r"width:\s*(\d+)", w)
                    add_image_from_src(img.get("src", ""), int(wm.group(1)) if wm else None)
                else:
                    add_paragraph_from_element(node)
                return
            for child in node.children:
                if getattr(child, "name", None):
                    walk(child)

        if body_nodes:
            walk(body_nodes[0])
        else:
            for block in soup.find_all(class_="editor-overlay-block"):
                walk(block)
            for el in soup.find_all(["p", "h1", "h2", "h3", "ul", "ol", "table"]):
                if el.find_parent(class_="doc-body"):
                    continue
                walk(el)

        if len(doc.paragraphs) == 0 and not doc.inline_shapes:
            doc.add_paragraph("Edited document")

        doc.save(output_path)
        return output_path

    elif task == "html-to-pdf":
        output_path = os.path.join(temp_dir, f"{uuid.uuid4()}_output.pdf")
        content = ""
        if files and len(files) > 0:
            with open(files[0], "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
        else:
            content = args.get("html_code", "")
            
        content = re.sub(r"<br\s*/?>", "\n", content)
        text = re.sub(r"<[^>]+>", "", content)
        c = canvas.Canvas(output_path, pagesize=letter)
        width, height = letter
        y = height - 50
        for line in text.split("\n"):
            line = line.strip()
            if not line: continue
            c.drawString(50, y, line[:100])
            y -= 15
            if y < 50:
                c.showPage()
                y = height - 50
        c.save()
        return output_path

    elif task == "summarize":
        input_path = files[0]
        ext = os.path.splitext(input_path)[1].lower()
        text = ""
        if ext == ".txt":
            with open(input_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        elif ext == ".pdf":
            with open(input_path, "rb") as f:
                pdf = PyPDF2.PdfReader(f)
                for page in pdf.pages:
                    text += page.extract_text() + "\n"
        elif ext == ".docx":
            text = extract_docx_text_safe(input_path)
        else:
            raise Exception("Unsupported file type. Use PDF, DOCX, or TXT.")

        if not text.strip():
            raise Exception("No readable text found in document. Try PDF or TXT, or re-save the Word file.")

        text = text[:12000]
        prompt = f"Summarize the following document clearly. Make a concise summary with important points that is easy to read.\n\nDOCUMENT:\n{text}"
        
        API_KEY = os.getenv("GROQ_API_KEY")
        if not API_KEY:
            raise Exception("AI summarizer is not configured. Set GROQ_API_KEY in your .env file.")
        client = Groq(api_key=API_KEY)
        completion = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3
        )
        return {"summary": completion.choices[0].message.content}

    elif task == "translate":
        input_path = files[0]
        from_lang = args.get("from_lang", "auto")
        to_lang = args.get("to_lang", "urdu")
        ext = os.path.splitext(input_path)[1].lower()
        text = ""
        if ext == ".pdf":
            doc = fitz.open(input_path)
            for page in doc:
                text += page.get_text()
            doc.close()
        elif ext == ".docx":
            text = extract_docx_text_safe(input_path)
        elif ext == ".txt":
            with open(input_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        else:
            raise Exception("Unsupported file type. Use PDF, DOCX, or TXT.")

        if not text.strip():
            raise Exception("No readable text found in document. Try PDF or TXT, or re-save the Word file.")

        lang_map = {
            "english": "en", "urdu": "ur", "arabic": "ar", "french": "fr",
            "spanish": "es", "german": "de", "hindi": "hi", "chinese": "zh-CN",
            "russian": "ru", "japanese": "ja", "turkish": "tr", "portuguese": "pt",
        }
        src = lang_map.get(from_lang, from_lang) if from_lang != "auto" else "auto"
        tgt = lang_map.get(to_lang, to_lang)

        translator = GoogleTranslator(source=src, target=tgt)
        chunks = [text[i:i+4500] for i in range(0, len(text), 4500)] if text else [""]
        translated_chunks = []
        for chunk in chunks:
            if chunk.strip():
                translated_chunks.append(translator.translate(chunk))
            else:
                translated_chunks.append("")
        
        return {
            "original_text": text[:5000],
            "translated_text": "\n".join(translated_chunks)
        }

    elif task == "ocr-tesseract":
        if not files:
            raise Exception("No file uploaded.")
        input_path = files[0]
        lang = args.get("lang", "eng")
        output_type = args.get("output_type", "json")
        psm = args.get("psm")
        if psm is not None and str(psm).strip() != "":
            psm = int(psm)
        else:
            psm = None

        result = ocr_file(input_path, lang=lang, temp_dir=temp_dir, psm=psm)

        if output_type == "txt":
            return write_text_output(result["text"], temp_dir)
        return result

    elif task == "ocr-status":
        return tesseract_status()

    elif task == "gmaps-scraper":
        industry = args.get("industry")
        location = args.get("location")
        max_results = int(args.get("max_results", 20))
        
        if not industry or not location:
            raise Exception("Industry and location are required.")
            
        try:
            import undetected_chromedriver as uc
        except ImportError:
            raise Exception("undetected_chromedriver is not installed.")

        output_path = os.path.join(temp_dir, f"leads_{uuid.uuid4().hex[:8]}.xlsx")

        class API_MapsScraper:
            def __init__(self):
                options = uc.ChromeOptions()
                options.add_argument("--headless=new")
                options.add_argument("--no-sandbox")
                options.add_argument("--disable-dev-shm-usage")
                options.add_argument("--user-agent=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/123.0.0.0 Safari/537.36")
                
                try:
                    self.driver = uc.Chrome(options=options, use_subprocess=True)
                except Exception as e:
                    chrome_options = Options()
                    chrome_options.add_argument("--headless=new")
                    chrome_options.add_argument("--no-sandbox")
                    self.driver = webdriver.Chrome(options=chrome_options)

            def search(self, industry, location, limit):
                query = f"{industry} in {location}"
                url = f"https://www.google.com/maps/search/{quote_plus(query)}"
                self.driver.get(url)
                time.sleep(5)
                
                try:
                    panel = self.driver.find_element(By.CSS_SELECTOR, '[role="feed"]')
                    for _ in range(5):
                        self.driver.execute_script("arguments[0].scrollTop += 1000;", panel)
                        time.sleep(2)
                except:
                    pass
                    
                anchors = self.driver.find_elements(By.CSS_SELECTOR, 'a[href*="/maps/place/"]')
                links = [a.get_attribute("href") for a in anchors if "/maps/place/" in (a.get_attribute("href") or "")]
                links = list(dict.fromkeys(links))[:limit]
                
                data = []
                for link in links:
                    try:
                        self.driver.get(link)
                        time.sleep(3)
                        soup = BeautifulSoup(self.driver.page_source, "html.parser")
                        
                        name = soup.select_one('h1.DUwDvf')
                        name = name.get_text(strip=True) if name else "N/A"
                        
                        address = soup.select_one('[data-item-id="address"]')
                        address = address.get_text(strip=True) if address else "N/A"
                        
                        phone = soup.select_one('[data-item-id^="phone"]')
                        phone = phone.get_text(strip=True) if phone else "N/A"
                        
                        website = soup.select_one('a[data-item-id="authority"]')
                        website = website.get("href") if website else "N/A"
                        
                        socials = {}
                        email = "N/A"
                        if website != "N/A":
                            try:
                                self.driver.get(website)
                                time.sleep(3)
                                web_text = self.driver.page_source
                                email_found = re.findall(r"[a-zA-Z0-9._%+\-]+@[a-zA-Z0-9.\-]+\.[a-zA-Z]{2,}", web_text)
                                email = email_found[0] if email_found else "N/A"
                                socials = extract_social_links(web_text)
                                self.driver.back()
                            except:
                                pass

                        data.append({
                            "Name": name,
                            "Address": address,
                            "Phone": phone,
                            "Website": website,
                            "Email": email,
                            "Facebook": socials.get("facebook", "N/A"),
                            "Instagram": socials.get("instagram", "N/A"),
                            "Twitter/X": socials.get("twitter", "N/A"),
                            "LinkedIn": socials.get("linkedin", "N/A"),
                            "YouTube": socials.get("youtube", "N/A"),
                            "Google Maps Link": link
                        })
                    except:
                        continue
                return data

            def close(self):
                self.driver.quit()

        scraper = API_MapsScraper()
        results = scraper.search(industry, location, max_results)
        scraper.close()
        
        if not results:
            raise Exception("No results found or search blocked.")
            
        df = pd.DataFrame(results)
        df.to_excel(output_path, index=False)
        return output_path

    else:
        raise Exception(f"Unknown task: {task}")


def main():
    try:
        input_data = sys.stdin.read()
        if not input_data.strip():
            print(json.dumps({"success": False, "error": "No input provided to worker"}))
            return
            
        data = json.loads(input_data)
        task = data.get("task")
        files = data.get("files", [])
        args = data.get("args", {})
        temp_dir = data.get("temp_dir", "temp")
        
        result = process_task(task, files, args, temp_dir)
        
        if isinstance(result, dict):
            print(json.dumps({"success": True, "data": result}))
        else:
            print(json.dumps({"success": True, "output": result}))
            
    except Exception as e:
        print(json.dumps({
            "success": False, 
            "error": str(e),
            "traceback": traceback.format_exc()
        }))

if __name__ == "__main__":
    main()
