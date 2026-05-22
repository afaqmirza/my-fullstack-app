from fastapi import FastAPI, File, UploadFile, Form, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
import shutil
import os
import zipfile
import re
import pdfplumber
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
from pdf2docx import Converter
from pypdf import PdfWriter, PdfReader
from PIL import Image, ImageOps
import img2pdf
import io
import fitz # PyMuPDF
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from typing import List, Optional
import time
import uuid
from urllib.parse import quote_plus
from bs4 import BeautifulSoup
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.common.by import By
from selenium.webdriver.support.ui import WebDriverWait
from selenium.webdriver.support import expected_conditions as EC
from groq import Groq
from docx import Document as DocxDocument
import PyPDF2
from deep_translator import GoogleTranslator
from docx_utils import extract_docx_text_safe

LANG_MAP = {
    "english": "en", "urdu": "ur", "arabic": "ar", "french": "fr",
    "spanish": "es", "german": "de", "hindi": "hi", "chinese": "zh-CN",
    "russian": "ru", "japanese": "ja", "turkish": "tr", "portuguese": "pt",
    "auto": "auto",
}
try:
    import pythoncom
    import win32com.client as win32
except ImportError:
    pythoncom = None
    win32 = None

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=["Content-Disposition", "X-Saved-Percent"]
)

IS_VERCEL = os.environ.get("VERCEL")
TEMP_DIR = "/tmp" if IS_VERCEL else "temp"
os.makedirs(TEMP_DIR, exist_ok=True)

def parse_page_ranges(range_str, total_pages):
    pages = set()
    parts = [p.strip() for p in range_str.split(",") if p.strip()]
    for part in parts:
        if "-" in part:
            try:
                start, end = part.split("-", 1)
                start = int(start.strip())
                end = int(end.strip())
                for p in range(start, end + 1):
                    if 1 <= p <= total_pages:
                        pages.add(p - 1)
            except: continue
        else:
            try:
                p = int(part.strip())
                if 1 <= p <= total_pages:
                    pages.add(p - 1)
            except: continue
    return sorted(pages)

@app.post("/api/convert/word-to-pdf")
async def convert_word_to_pdf(file: UploadFile = File(...)):
    if not file.filename.endswith((".doc", ".docx")):
        return {"error": "Invalid file format. Please upload a Word document."}
        
    input_path = os.path.abspath(f"temp/{file.filename}")
    output_filename = f"{file.filename.rsplit('.', 1)[0]}.pdf"
    output_path = os.path.abspath(f"temp/{output_filename}")
    
    # Save the uploaded file to disk
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    if pythoncom is None or win32 is None:
        return {"error": "pywin32 is required for Word conversion on Windows."}

    try:
        pythoncom.CoInitialize()
        word = win32.Dispatch("Word.Application")
        word.Visible = False
        word.DisplayAlerts = 0
        
        doc = word.Documents.Open(input_path)
        doc.SaveAs(output_path, FileFormat=17) # 17 is the code for PDF
        doc.Close(False)
        word.Quit()
        pythoncom.CoUninitialize()
        
        return FileResponse(
            output_path, 
            media_type="application/pdf", 
            filename=output_filename,
            headers={"Content-Disposition": f"attachment; filename={output_filename}"}
        )
    except Exception as e:
        if 'word' in locals() and word:
            word.Quit()
        return {"error": f"Word conversion error: {str(e)}"}

@app.post("/api/convert/pdf-to-ppt")
async def convert_pdf_to_ppt(file: UploadFile = File(...)):
    if not file.filename.lower().endswith(".pdf"):
        return {"error": "Invalid file format. Please upload a PDF document."}
        
    input_path = f"temp/{file.filename}"
    output_filename = f"{file.filename.rsplit('.', 1)[0]}.pptx"
    output_path = f"temp/{output_filename}"
    
    # Save the uploaded file to disk
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    try:
        # Create a new presentation
        prs = Presentation()
        
        with pdfplumber.open(input_path) as pdf:
            for page in pdf.pages:
                # Create a slide
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # Convert page to image and add to slide
                # This is the "integrated method" similar to local_convertor.py
                img_path = f"temp/page_{page.page_number}.png"
                page.to_image(resolution=150).save(img_path)
                
                # Add image to slide, filling the whole slide
                slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
                
                # Clean up temp image
                if os.path.exists(img_path):
                    os.remove(img_path)
        
        # Save as PPTX
        prs.save(output_path)
        
        return FileResponse(
            output_path, 
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", 
            filename=output_filename,
            headers={"Content-Disposition": f"attachment; filename={output_filename}"}
        )
    except Exception as e:
        return {"error": str(e)}

@app.post("/api/convert/pdf-to-image")
async def convert_pdf_to_image(file: UploadFile = File(...)):
    if not file.filename.lower().endswith(".pdf"):
        return {"error": "Invalid file format. Please upload a PDF document."}
        
    input_path = f"temp/{file.filename}"
    base_name = file.filename.rsplit('.', 1)[0]
    output_filename = f"{base_name}_images.zip"
    output_path = f"temp/{output_filename}"
    
    # Save the uploaded file to disk
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    try:
        image_dir = f"temp/{base_name}_images"
        os.makedirs(image_dir, exist_ok=True)
        
        with pdfplumber.open(input_path) as pdf:
            for page in pdf.pages:
                image_path = f"{image_dir}/page_{page.page_number}.png"
                page.to_image(resolution=300).save(image_path)
                
        # Create ZIP
        with zipfile.ZipFile(output_path, 'w') as zipf:
            for root, _, files in os.walk(image_dir):
                for f in files:
                    zipf.write(os.path.join(root, f), f)
                    
        return FileResponse(
            output_path, 
            media_type="application/zip", 
            filename=output_filename,
            headers={"Content-Disposition": f"attachment; filename={output_filename}"}
        )
    except Exception as e:
        return {"error": str(e)}

@app.post("/api/convert/pdf-to-excel")
async def convert_pdf_to_excel(file: UploadFile = File(...)):
    if not file.filename.lower().endswith(".pdf"):
        return {"error": "Invalid file format. Please upload a PDF document."}
        
    input_path = f"temp/{file.filename}"
    output_filename = f"{file.filename.rsplit('.', 1)[0]}.xlsx"
    output_path = f"temp/{output_filename}"
    
    # Save the uploaded file to disk
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    try:
        all_text = ""
        with pdfplumber.open(input_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    all_text += " " + text

        # New logic: Extract sentences instead of tables
        clean_text = " ".join(all_text.split())
        sentences = re.split(r'(?<=[.!?])\s+(?=[A-Z])', clean_text)
        sentences = [s.strip() for s in sentences if s.strip()]

        if sentences:
            df = pd.DataFrame(sentences, columns=["Sentence"])
            with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
                df.to_excel(writer, index=False)
                # Apply column width logic from the new tool
                ws = writer.sheets["Sheet1"]
                ws.column_dimensions["A"].width = 120
        else:
            return {"error": "No content found in PDF to convert."}
        
        return FileResponse(
            output_path, 
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", 
            filename=output_filename,
            headers={"Content-Disposition": f"attachment; filename={output_filename}"}
        )
    except Exception as e:
        return {"error": str(e)}

@app.post("/api/convert/pdf-to-word")
async def convert_pdf_to_word(file: UploadFile = File(...)):
    if not file.filename.lower().endswith(".pdf"):
        return {"error": "Invalid file format. Please upload a PDF document."}
        
    input_path = f"temp/{file.filename}"
    output_filename = f"{file.filename.rsplit('.', 1)[0]}.docx"
    output_path = f"temp/{output_filename}"
    
    # Save the uploaded file to disk
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    try:
        # Conversion logic from the new tool
        cv = Converter(input_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()
        
        return FileResponse(
            output_path, 
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", 
            filename=output_filename,
            headers={"Content-Disposition": f"attachment; filename={output_filename}"}
        )
    except Exception as e:
        return {"error": str(e)}

@app.post("/api/convert/pdf-merge")
async def merge_pdfs(files: List[UploadFile] = File(...)):
    if len(files) < 2:
        return {"error": "Please upload at least 2 PDF files to merge."}
        
    output_filename = "merged_document.pdf"
    output_path = f"temp/{output_filename}"
    
    try:
        writer = PdfWriter()
        temp_files = []
        
        for file in files:
            if not file.filename.lower().endswith(".pdf"):
                return {"error": f"File {file.filename} is not a PDF."}
                
            temp_path = f"temp/{file.filename}"
            with open(temp_path, "wb") as buffer:
                shutil.copyfileobj(file.file, buffer)
            temp_files.append(temp_path)
            
            reader = PdfReader(temp_path)
            for page in reader.pages:
                writer.add_page(page)
        
        with open(output_path, "wb") as out_file:
            writer.write(out_file)
            
        # Clean up temp input files
        for path in temp_files:
            if os.path.exists(path):
                os.remove(path)
                
        return FileResponse(
            output_path, 
            media_type="application/pdf", 
            filename=output_filename,
            headers={"Content-Disposition": f"attachment; filename={output_filename}"}
        )
    except Exception as e:
        return {"error": str(e)}

@app.post("/api/convert/pdf-split")
async def split_pdf(
    file: UploadFile = File(...),
    mode: str = "range", # "range", "each", "parts"
    range_str: str = None,
    parts_count: int = None
):
    if not file.filename.lower().endswith(".pdf"):
        return {"error": "Invalid file format. Please upload a PDF document."}
        
    input_path = f"temp/{file.filename}"
    base_name = file.filename.rsplit('.', 1)[0]
    
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    try:
        reader = PdfReader(input_path)
        total_pages = len(reader.pages)
        output_paths = []
        
        if mode == "range":
            if not range_str:
                return {"error": "Page range is required for range mode."}
            pages = parse_page_ranges(range_str, total_pages)
            if not pages:
                return {"error": "No valid pages found in range."}
            
            writer = PdfWriter()
            for idx in pages:
                writer.add_page(reader.pages[idx])
            output_filename = f"{base_name}_split.pdf"
            output_path = f"temp/{output_filename}"
            with open(output_path, "wb") as f:
                writer.write(f)
            return FileResponse(output_path, media_type="application/pdf", filename=output_filename)
            
        elif mode == "each" or mode == "parts":
            # These modes return a ZIP
            output_zip_filename = f"{base_name}_split_parts.zip"
            output_zip_path = f"temp/{output_zip_filename}"
            image_dir = f"temp/{base_name}_split"
            os.makedirs(image_dir, exist_ok=True)
            
            if mode == "each":
                for i in range(total_pages):
                    writer = PdfWriter()
                    writer.add_page(reader.pages[i])
                    out_part = f"{image_dir}/{base_name}_page_{i+1}.pdf"
                    with open(out_part, "wb") as f:
                        writer.write(f)
            else: # parts
                if not parts_count or parts_count < 2:
                    return {"error": "Valid number of parts (>=2) is required."}
                pages_per_part = total_pages // parts_count
                remainder = total_pages % parts_count
                start = 0
                for part in range(parts_count):
                    extra = 1 if part < remainder else 0
                    end = start + pages_per_part + extra
                    if start == end: continue
                    writer = PdfWriter()
                    for idx in range(start, end):
                        writer.add_page(reader.pages[idx])
                    out_part = f"{image_dir}/{base_name}_part_{part+1}.pdf"
                    with open(out_part, "wb") as f:
                        writer.write(f)
                    start = end
            
            with zipfile.ZipFile(output_zip_path, 'w') as zipf:
                for root, _, files in os.walk(image_dir):
                    for f in files:
                        zipf.write(os.path.join(root, f), f)
            return FileResponse(output_zip_path, media_type="application/zip", filename=output_zip_filename)
            
    except Exception as e:
        return {"error": str(e)}

@app.post("/api/convert/pdf-compress")
async def api_compress_pdf(file: UploadFile = File(...)):
    if not file.filename.lower().endswith(".pdf"):
        return {"error": "Invalid file format. Please upload a PDF document."}
        
    input_path = f"temp/{file.filename}"
    output_filename = f"{file.filename.rsplit('.', 1)[0]}_compressed.pdf"
    output_path = f"temp/{output_filename}"
    
    # Save the uploaded file to disk
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    try:
        reader = PdfReader(input_path)
        writer = PdfWriter()

        for page in reader.pages:
            # Basic compression in pypdf
            page.compress_content_streams()
            writer.add_page(page)

        # Remove duplicate objects
        writer.compress_identical_objects(remove_identicals=True, remove_orphans=True)

        with open(output_path, "wb") as f:
            writer.write(f)
            
        orig_size = os.path.getsize(input_path)
        new_size = os.path.getsize(output_path)
        saved_bytes = orig_size - new_size
        saved_pct = (saved_bytes / orig_size * 100) if orig_size > 0 else 0
        
        return FileResponse(
            output_path, 
            media_type="application/pdf", 
            filename=output_filename,
            headers={
                "Content-Disposition": f"attachment; filename={output_filename}",
                "X-Saved-Percent": f"{saved_pct:.1f}"
            }
        )
    except Exception as e:
        return {"error": str(e)}

@app.post("/api/convert/image-to-pdf")
async def api_image_to_pdf(files: List[UploadFile] = File(...)):
    output_filename = "images_converted.pdf"
    output_path = f"temp/{output_filename}"
    
    image_paths = []
    for file in files:
        file_path = f"temp/{file.filename}"
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        # Prepare image (remove alpha, rotate if needed)
        try:
            img = Image.open(file_path)
            img = ImageOps.exif_transpose(img)
            if img.mode in ("RGBA", "LA", "P"):
                bg = Image.new("RGB", img.size, (255, 255, 255))
                if img.mode == "P":
                    img = img.convert("RGBA")
                bg.paste(img, mask=img.split()[-1] if img.mode == "RGBA" else None)
                img = bg
            elif img.mode != "RGB":
                img = img.convert("RGB")
            
            # Save to a temporary buffer as JPEG to ensure compatibility
            buf_path = f"temp/processed_{file.filename}.jpg"
            img.save(buf_path, format="JPEG", quality=95)
            image_paths.append(buf_path)
        except Exception as e:
            print(f"Error processing {file.filename}: {e}")
            continue

    if not image_paths:
        return {"error": "No valid images processed."}

    try:
        pdf_bytes = img2pdf.convert(image_paths)
        with open(output_path, "wb") as f:
            f.write(pdf_bytes)
        
        return FileResponse(output_path, media_type="application/pdf", filename=output_filename)
    except Exception as e:
        return {"error": str(e)}

@app.post("/api/convert/pdf-protect")
async def api_pdf_protect(
    file: UploadFile = File(...),
    password: str = Form(...),
    permissions: Optional[str] = Form(None) # Optional JSON string of perms
):
    input_path = f"temp/{file.filename}"
    output_filename = f"protected_{file.filename}"
    output_path = f"temp/{output_filename}"
    
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    try:
        reader = PdfReader(input_path)
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
            
        writer.encrypt(user_password=password, owner_password=password, use_128bit=True)
        
        with open(output_path, "wb") as f:
            writer.write(f)
            
        return FileResponse(output_path, media_type="application/pdf", filename=output_filename)
    except Exception as e:
        return {"error": str(e)}

@app.post("/api/convert/pdf-unlock")
async def api_pdf_unlock(
    file: UploadFile = File(...),
    password: str = Form(...)
):
    input_path = f"temp/{file.filename}"
    output_filename = f"unlocked_{file.filename}"
    output_path = f"temp/{output_filename}"
    
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    try:
        reader = PdfReader(input_path)
        if reader.is_encrypted:
            result = reader.decrypt(password)
            if result == 0:
                return {"error": "Incorrect password."}
        else:
            return {"error": "PDF is not encrypted."}
            
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
            
        with open(output_path, "wb") as f:
            writer.write(f)
            
        return FileResponse(output_path, media_type="application/pdf", filename=output_filename)
    except Exception as e:
        return {"error": str(e)}

@app.post("/api/convert/pdf-rotate")
async def api_pdf_rotate(
    file: UploadFile = File(...),
    angle: int = Form(...),
    pages: str = Form("ALL")
):
    input_path = f"temp/{file.filename}"
    output_filename = f"rotated_{file.filename}"
    output_path = f"temp/{output_filename}"
    
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    try:
        reader = PdfReader(input_path)
        writer = PdfWriter()
        total_pages = len(reader.pages)
        
        target_pages = parse_page_ranges(pages, total_pages) if pages.upper() != "ALL" else range(total_pages)
        
        for i, page in enumerate(reader.pages):
            if i in target_pages:
                page.rotate(angle)
            writer.add_page(page)
            
        with open(output_path, "wb") as f:
            writer.write(f)
            
        return FileResponse(output_path, media_type="application/pdf", filename=output_filename)
    except Exception as e:
        return {"error": str(e)}

@app.post("/api/convert/pdf-watermark")
async def api_pdf_watermark(
    file: UploadFile = File(...),
    text: str = Form(...),
    opacity: float = Form(0.3)
):
    input_path = f"temp/{file.filename}"
    output_filename = f"watermarked_{file.filename}"
    output_path = f"temp/{output_filename}"
    
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    try:
        # Create watermark overlay
        packet = io.BytesIO()
        can = canvas.Canvas(packet, pagesize=letter)
        can.setFillColorRGB(0, 0, 0, alpha=opacity)
        can.setFont("Helvetica", 40)
        can.saveState()
        can.translate(200, 400)
        can.rotate(45)
        can.drawString(0, 0, text)
        can.restoreState()
        can.save()
        packet.seek(0)
        
        watermark_reader = PdfReader(packet)
        wm_page = watermark_reader.pages[0]
        
        reader = PdfReader(input_path)
        writer = PdfWriter()
        
        for page in reader.pages:
            page.merge_page(wm_page)
            writer.add_page(page)
            
        with open(output_path, "wb") as f:
            writer.write(f)
            
        return FileResponse(output_path, media_type="application/pdf", filename=output_filename)
    except Exception as e:
        return {"error": str(e)}

@app.post("/api/convert/pdf-edit")
async def api_pdf_edit(
    file: UploadFile = File(...),
    edit_text: Optional[str] = Form(None),
    edit_image: Optional[UploadFile] = File(None),
    x: int = Form(100),
    y: int = Form(100)
):
    input_path = f"temp/{file.filename}"
    output_filename = f"edited_{file.filename}"
    output_path = f"temp/{output_filename}"
    
    with open(input_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    try:
        doc = fitz.open(input_path)
        page = doc[0] # Edit first page for now
        
        if edit_text:
            page.insert_text((x, y), edit_text, fontsize=14, color=(0, 0, 0))
            
        if edit_image:
            img_path = f"temp/edit_img_{edit_image.filename}"
            with open(img_path, "wb") as buffer:
                shutil.copyfileobj(edit_image.file, buffer)
            rect = fitz.Rect(x, y, x + 150, y + 150)
            page.insert_image(rect, filename=img_path)
            
        doc.save(output_path)
        return FileResponse(output_path, media_type="application/pdf", filename=output_filename)
    except Exception as e:
        return {"error": str(e)}

@app.post("/api/convert/html-to-pdf")
async def api_html_to_pdf(
    file: Optional[UploadFile] = File(None),
    html_code: Optional[str] = Form(None)
):
    output_filename = "html_converted.pdf"
    output_path = f"temp/{output_filename}"
    
    content = ""
    if file:
        file_content = await file.read()
        content = file_content.decode("utf-8", errors="ignore")
    elif html_code:
        content = html_code
        
    if not content:
        return {"error": "No HTML content provided."}
        
    try:
        # Simple tag stripping as in the provided script
        import re
        content = re.sub(r"<br\s*/?>", "\n", content)
        text = re.sub(r"<[^>]+>", "", content)
        
        c = canvas.Canvas(output_path, pagesize=letter)
        width, height = letter
        y = height - 50
        
        for line in text.split("\n"):
            line = line.strip()
            if not line: continue
            c.drawString(50, y, line[:100])
            y -= 15
            if y < 50:
                c.showPage()
                y = height - 50
        
        c.save()
        return FileResponse(output_path, media_type="application/pdf", filename=output_filename)
    except Exception as e:
        return {"error": str(e)}

# ─── Automation Tools Helpers ─────────────────────────────────────────────

SOCIAL_PATTERNS = {
    "facebook":  re.compile(r"https?://(www\.)?facebook\.com/(?!sharer|share|dialog|plugins|tr\b)[^\s\"'<>?&]+"),
    "instagram": re.compile(r"https?://(www\.)?instagram\.com/[^\s\"'<>?&]+"),
    "twitter":   re.compile(r"https?://(www\.)?(twitter\.com|x\.com)/[^\s\"'<>?&]+"),
    "linkedin":  re.compile(r"https?://(www\.)?linkedin\.com/(company|in)/[^\s\"'<>?&]+"),
    "youtube":   re.compile(r"https?://(www\.)?youtube\.com/(channel|c|user|@)[^\s\"'<>?&]+"),
}

def extract_social_links(text: str) -> dict:
    result = {}
    for platform, pat in SOCIAL_PATTERNS.items():
        matches = [m.group(0) for m in pat.finditer(text)]
        seen = set()
        unique = []
        for m in matches:
            base = m.rstrip("/")
            if base not in seen:
                seen.add(base)
                unique.append(base)
        if unique:
            result[platform] = unique[0]
    return result

# ─── Automation Tools: Google Maps Scraper ──────────────────────────────────

@app.post("/api/automation/gmaps-scraper")
async def api_gmaps_scraper(
    industry: str = Form(...),
    location: str = Form(...),
    max_results: int = Form(20)
):
    try:
        import undetected_chromedriver as uc
    except ImportError:
        return JSONResponse({"error": "undetected_chromedriver not installed"}, status_code=500)

    output_filename = f"leads_{uuid.uuid4().hex[:8]}.xlsx"
    output_path = f"temp/{output_filename}"

    # Minimal Scraper Implementation for API
    class API_MapsScraper:
        def __init__(self):
            options = uc.ChromeOptions()
            options.add_argument("--headless=new")
            options.add_argument("--no-sandbox")
            options.add_argument("--disable-dev-shm-usage")
            # Set a common user agent
            options.add_argument("--user-agent=Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/123.0.0.0 Safari/537.36")
            
            try:
                self.driver = uc.Chrome(options=options, use_subprocess=True)
            except Exception as e:
                # Fallback to standard selenium if uc fails
                chrome_options = Options()
                chrome_options.add_argument("--headless=new")
                chrome_options.add_argument("--no-sandbox")
                self.driver = webdriver.Chrome(options=chrome_options)

        def search(self, industry, location, limit):
            query = f"{industry} in {location}"
            url = f"https://www.google.com/maps/search/{quote_plus(query)}"
            self.driver.get(url)
            time.sleep(5) # Wait for initial load
            
            # Simple scrolling to collect results
            results = []
            try:
                panel = self.driver.find_element(By.CSS_SELECTOR, '[role="feed"]')
                for _ in range(5):
                    self.driver.execute_script("arguments[0].scrollTop += 1000;", panel)
                    time.sleep(2)
            except:
                pass
                
            anchors = self.driver.find_elements(By.CSS_SELECTOR, 'a[href*="/maps/place/"]')
            links = [a.get_attribute("href") for a in anchors if "/maps/place/" in (a.get_attribute("href") or "")]
            links = list(dict.fromkeys(links))[:limit]
            
            data = []
            for link in links:
                try:
                    self.driver.get(link)
                    time.sleep(3)
                    soup = BeautifulSoup(self.driver.page_source, "html.parser")
                    
                    name = soup.select_one('h1.DUwDvf')
                    name = name.get_text(strip=True) if name else "N/A"
                    
                    address = soup.select_one('[data-item-id="address"]')
                    address = address.get_text(strip=True) if address else "N/A"
                    
                    phone = soup.select_one('[data-item-id^="phone"]')
                    phone = phone.get_text(strip=True) if phone else "N/A"
                    
                    website = soup.select_one('a[data-item-id="authority"]')
                    website = website.get("href") if website else "N/A"
                    
                    socials = {}
                    email = "N/A"
                    if website != "N/A":
                        try:
                            self.driver.get(website)
                            time.sleep(3)
                            web_soup = BeautifulSoup(self.driver.page_source, "html.parser")
                            web_text = self.driver.page_source
                            
                            # Extract Email
                            email_found = re.findall(r"[a-zA-Z0-9._%+\-]+@[a-zA-Z0-9.\-]+\.[a-zA-Z]{2,}", web_text)
                            email = email_found[0] if email_found else "N/A"
                            
                            # Extract Socials
                            socials = extract_social_links(web_text)
                            self.driver.back()
                        except:
                            pass

                    data.append({
                        "Name": name,
                        "Address": address,
                        "Phone": phone,
                        "Website": website,
                        "Email": email,
                        "Facebook": socials.get("facebook", "N/A"),
                        "Instagram": socials.get("instagram", "N/A"),
                        "Twitter/X": socials.get("twitter", "N/A"),
                        "LinkedIn": socials.get("linkedin", "N/A"),
                        "YouTube": socials.get("youtube", "N/A"),
                        "Google Maps Link": link
                    })
                except:
                    continue
            return data

        def close(self):
            self.driver.quit()

    try:
        scraper = API_MapsScraper()
        results = scraper.search(industry, location, max_results)
        scraper.close()
        
        if not results:
            return JSONResponse({"error": "No results found or search blocked."}, status_code=404)
            
        df = pd.DataFrame(results)
        df.to_excel(output_path, index=False)
        
        return FileResponse(output_path, media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", filename="business_leads.xlsx")
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)

# ─── Intelligence Tools: AI Summarizer ──────────────────────────────────────
# =========================
# CONFIG
# =========================
API_KEY = os.getenv("GROQ_API_KEY")

client = Groq(api_key=API_KEY)

@app.post("/api/intelligence/summarize")
async def api_summarize(file: UploadFile = File(...)):
    temp_path = f"{TEMP_DIR}/{uuid.uuid4()}_{file.filename}"
    with open(temp_path, "wb") as f:
        shutil.copyfileobj(file.file, f)
    
    try:
        text = ""
        ext = os.path.splitext(file.filename)[1].lower()
        
        if ext == ".txt":
            with open(temp_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        elif ext == ".pdf":
            with open(temp_path, "rb") as f:
                pdf = PyPDF2.PdfReader(f)
                for page in pdf.pages:
                    text += page.extract_text() + "\n"
        elif ext == ".docx":
            text = extract_docx_text_safe(temp_path)
        else:
            return JSONResponse({"error": "Unsupported file type. Use PDF, DOCX, or TXT."}, status_code=400)
        
        if not text.strip():
            return JSONResponse({"error": "No readable text found in document"}, status_code=400)

        if not API_KEY:
            return JSONResponse({"error": "AI summarizer is not configured. Set GROQ_API_KEY in .env"}, status_code=500)
            
        # Truncate text to stay within token limits if necessary
        text = text[:12000]
        
        prompt = f"Summarize the following document clearly. Make a concise summary with important points that is easy to read.\n\nDOCUMENT:\n{text}"
        
        completion = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3
        )
        
        summary = completion.choices[0].message.content
        return {"summary": summary}
        
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)

@app.post("/api/intelligence/translate")
async def api_translate(
    file: UploadFile = File(...),
    from_lang: str = Form("auto"),
    to_lang: str = Form("urdu")
):
    temp_path = f"{TEMP_DIR}/{uuid.uuid4()}_{file.filename}"
    with open(temp_path, "wb") as f:
        shutil.copyfileobj(file.file, f)
    
    try:
        text = ""
        ext = os.path.splitext(file.filename)[1].lower()
        
        # Extraction logic
        if ext == ".pdf":
            doc = fitz.open(temp_path)
            for page in doc:
                text += page.get_text()
            doc.close()
        elif ext == ".docx":
            text = extract_docx_text_safe(temp_path)
        elif ext == ".txt":
            with open(temp_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        else:
            return JSONResponse({"error": "Unsupported file type. Use PDF, DOCX, or TXT."}, status_code=400)
        
        if not text.strip():
            return JSONResponse({"error": "No readable text found"}, status_code=400)

        src = LANG_MAP.get(from_lang, from_lang) if from_lang != "auto" else "auto"
        tgt = LANG_MAP.get(to_lang, to_lang)
        translator = GoogleTranslator(source=src, target=tgt)
        
        # Split text into chunks for translator (limit is 5000 chars)
        chunks = [text[i:i+4500] for i in range(0, len(text), 4500)] if text else [""]
        translated_chunks = []
        
        for chunk in chunks:
            if chunk.strip():
                translated_chunks.append(translator.translate(chunk))
            else:
                translated_chunks.append("")
            
        translated_text = "\n".join(translated_chunks)
        
        return {
            "original_text": text[:5000], # Return snippet for preview
            "translated_text": translated_text
        }
        
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)

